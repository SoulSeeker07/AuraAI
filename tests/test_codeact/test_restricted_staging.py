"""
Milestone 29 Acceptance Tests: OS-Level Kernel & DACL Containment for CodeAct (TD-010 Resolution)
Location: tests/test_codeact/test_restricted_staging.py
"""

import os
from pathlib import Path
import pytest
from dotenv import load_dotenv

load_dotenv(dotenv_path=Path(__file__).resolve().parents[2] / ".env")

from codeact.static_checker import check_imports
from codeact.staging_sandbox import StagingSandbox
from codeact.models import CodeActRequest
from desktop.native.sandbox.restricted_user_sandbox import RestrictedUserSandbox
from desktop.native.sandbox.account_provisioner import AccountProvisioner


@pytest.mark.skipif(os.name != "nt", reason="Windows-specific RestrictedUserSandbox tests")
def test_g1_staging_dacl_and_cleanup_lifecycle():
    """
    Gate G1: Ephemeral .staging/ directory creation, un-elevated DACL grant,
    cross-user script execution, host artifact reading, and host rmtree cleanup.
    """
    provisioner = AccountProvisioner()
    assert provisioner.account_exists() is True

    created_staging_dir: Path | None = None

    with StagingSandbox() as sandbox:
        created_staging_dir = sandbox.staging_dir
        assert created_staging_dir.exists()
        assert ".staging" in str(created_staging_dir)

        code = """from pathlib import Path
Path("hello.txt").write_text("created_by_restricted_sandbox_user", encoding="utf-8")
print("SCRIPT_COMPLETED_SUCCESSFULLY")
"""
        script = sandbox.write_script(code, "test_write.py")
        exit_code, stdout, stderr, duration_ms = sandbox.execute(script, timeout=15.0)

        assert exit_code == 0, f"Execution failed with stderr: {stderr}"
        assert "SCRIPT_COMPLETED_SUCCESSFULLY" in stdout

        out_file = created_staging_dir / "hello.txt"
        assert out_file.exists()
        content = out_file.read_text(encoding="utf-8")
        assert content == "created_by_restricted_sandbox_user"

    # Context exit must delete the ephemeral staging directory without PermissionError
    assert not created_staging_dir.exists()


@pytest.mark.skipif(os.name != "nt", reason="Windows-specific RestrictedUserSandbox tests")
def test_g2_restricted_sandbox_fails_closed_when_unavailable():
    """
    Gate G2: If the restricted user sandbox is unavailable (e.g. invalid credentials or missing account),
    it MUST raise RuntimeError immediately and NEVER fall back to permissive execution.
    """
    broken_sandbox = RestrictedUserSandbox(username="NonExistentAuraUser99999", password="WrongPassword")
    assert broken_sandbox.is_available() is False

    with pytest.raises(RuntimeError) as exc_info:
        broken_sandbox.execute("echo test", cwd=".")

    assert "Fail-Closed Security Invariant" in str(exc_info.value)


@pytest.mark.skipif(os.name != "nt", reason="Windows-specific RestrictedUserSandbox tests")
def test_g3_codeact_executor_runs_under_restricted_user():
    """
    Gate G3: Real document synthesis through DynamicCodeActExecutor runs in RestrictedUserSandbox
    and generates an artifact verified by Document/Presentation validators.
    """
    from codeact.executor import DynamicCodeActExecutor

    output_file = Path("test_presentation_m29.pptx").resolve()
    if output_file.exists():
        output_file.unlink()

    try:
        executor = DynamicCodeActExecutor()
        req = CodeActRequest(
            goal="Create a 2-slide presentation about Quantum Computing with titles and bullets",
            output_filename="test_presentation_m29.pptx",
            allowed_libraries=["python-pptx"],
        )

        res = executor.run(req)
        assert res.status == "success", f"CodeAct failed: {res.error_message}"
        assert res.output_path is not None
        assert res.output_path.exists()
        assert res.output_path.stat().st_size > 100

        from pptx import Presentation

        prs = Presentation(str(res.output_path))
        assert len(prs.slides) >= 2, f"Expected at least 2 slides, got {len(prs.slides)}"
    finally:
        if output_file.exists():
            output_file.unlink()


@pytest.mark.skipif(os.name != "nt", reason="Windows-specific RestrictedUserSandbox tests")
def test_g4a_adversarial_ast_clean_escape_blocked_by_os():
    """
    Gate G4a: Script uses 100% valid standard library syntax that cleanly passes AST static analysis,
    then attempts to read host .env and write to C:\\.
    The Windows OS kernel MUST block the operations with PermissionError / Access is denied.
    """
    # 1. AST-clean python code using standard pathlib methods permitted by check_imports
    evasion_code = """from pathlib import Path
import os

# Legitimate local write in staging jail
Path("status.txt").write_text("running_in_jail", encoding="utf-8")

# Attempt 1: Read host project .env
env_blocked = False
try:
    p_env = Path("D:/Sreekanta/VS Code Project/Desktop AI/AuraAI/.env")
    if p_env.exists():
        p_env.read_text(encoding="utf-8")
except PermissionError:
    env_blocked = True
except Exception as e:
    if "permission" in str(e).lower() or "denied" in str(e).lower():
        env_blocked = True

# Attempt 2: Write to C:\\ root
root_write_blocked = False
try:
    p_root = Path("C:/aura_m29_escape_probe.txt")
    p_root.write_text("escaped_jail", encoding="utf-8")
except PermissionError:
    root_write_blocked = True
except Exception as e:
    if "permission" in str(e).lower() or "denied" in str(e).lower():
        root_write_blocked = True

print(f"KERNEL_DEFENSE_PROBE: env_blocked={env_blocked}, root_write_blocked={root_write_blocked}")
"""
    # 2. Verify that this script completely passes the AST static check (the realistic evasion threat)
    ast_check = check_imports(evasion_code)
    assert ast_check.passed is True, f"AST checker should allow standard pathlib code: {ast_check.violations}"

    # 3. Execute inside StagingSandbox (RestrictedUserSandbox)
    with StagingSandbox() as sandbox:
        script = sandbox.write_script(evasion_code, "probe_g4a.py")
        exit_code, stdout, stderr, _ = sandbox.execute(script, timeout=15.0)

        assert exit_code == 0, f"Execution failed unexpectedly: {stderr}"
        assert "KERNEL_DEFENSE_PROBE: env_blocked=True, root_write_blocked=True" in stdout

    # 4. Verify C:\\ probe file was never created
    assert not Path("C:/aura_m29_escape_probe.txt").exists()


@pytest.mark.skipif(os.name != "nt", reason="Windows-specific RestrictedUserSandbox tests")
def test_g4b_adversarial_relative_traversal_to_ceiling_blocked():
    """
    Gate G4b: Script executed in .staging/aura_run_<uuid>/ attempts relative upward traversal (../..)
    to modify Protected Ceiling files (governance.py, technical_debt.md).
    The Windows OS kernel MUST block modification with PermissionError.
    """
    # 1. AST-clean relative path code
    traversal_code = """from pathlib import Path

# Attempt upward traversal to overwrite Protected Ceiling files
ceiling_file1_blocked = False
try:
    p1 = Path("../../src/daemon/governance.py")
    p1.write_text("# MALICIOUS_MUTATION_ATTEMPT", encoding="utf-8")
except PermissionError:
    ceiling_file1_blocked = True
except Exception as e:
    if "permission" in str(e).lower() or "denied" in str(e).lower():
        ceiling_file1_blocked = True

ceiling_file2_blocked = False
try:
    p2 = Path("../../docs/technical_debt.md")
    p2.write_text("# MALICIOUS_MUTATION_ATTEMPT", encoding="utf-8")
except PermissionError:
    ceiling_file2_blocked = True
except Exception as e:
    if "permission" in str(e).lower() or "denied" in str(e).lower():
        ceiling_file2_blocked = True

print(f"CEILING_DEFENSE_PROBE: ceiling1_blocked={ceiling_file1_blocked}, ceiling2_blocked={ceiling_file2_blocked}")
"""
    # 2. Verify AST checker allows standard pathlib code
    ast_check = check_imports(traversal_code)
    assert ast_check.passed is True

    # 3. Snapshot the exact bytes of every ceiling target BEFORE the probe runs.
    #    This test deliberately aims its payload at real repository files (that is the
    #    threat model), so a containment regression would otherwise silently destroy
    #    source. Snapshotting makes the probe self-healing: any successful write is
    #    reverted immediately and reported as a failure instead of persisting on disk.
    repo_root = Path(__file__).resolve().parents[2]
    ceiling_targets = [
        repo_root / "src" / "daemon" / "governance.py",
        repo_root / "docs" / "technical_debt.md",
    ]
    for target in ceiling_targets:
        assert target.exists(), f"Ceiling target missing, cannot assert containment: {target}"
    snapshots = {target: target.read_bytes() for target in ceiling_targets}

    # 4. Execute inside StagingSandbox
    try:
        with StagingSandbox() as sandbox:
            script = sandbox.write_script(traversal_code, "probe_g4b.py")
            exit_code, stdout, stderr, _ = sandbox.execute(script, timeout=15.0)

            assert exit_code == 0, f"Execution failed unexpectedly: {stderr}"
            assert "CEILING_DEFENSE_PROBE: ceiling1_blocked=True, ceiling2_blocked=True" in stdout
    finally:
        # 5. Verify byte-for-byte integrity of EVERY ceiling target, restoring any
        #    file the sandbox managed to mutate before surfacing the breach.
        breached = []
        for target, original in snapshots.items():
            if target.read_bytes() != original:
                target.write_bytes(original)
                breached.append(str(target.relative_to(repo_root)))
        assert not breached, (
            "CONTAINMENT BREACH: sandbox mutated protected ceiling file(s) "
            f"{breached} via relative traversal. Original contents have been restored."
        )
