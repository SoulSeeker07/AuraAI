"""
Unit Tests for CodeAct Output Validators
Location: tests/test_codeact/test_validators.py
"""

from pathlib import Path
import pytest
from codeact.models import CodeActRequest, ExecutionAttempt
from codeact.validators import validate


def test_validate_text_success(tmp_path):
    out_file = tmp_path / "summary.md"
    out_file.write_text("# Title\n\nThis is a non-empty summary content.\n" * 5, encoding="utf-8")

    req = CodeActRequest(goal="summary", output_filename="summary.md")
    attempt = ExecutionAttempt(
        attempt_number=1,
        code="print('done')",
        stdout="done",
        stderr="",
        traceback=None,
        exit_code=0,
        duration_ms=100,
    )

    res = validate(req, tmp_path, attempt)
    assert res.passed is True
    assert dict(res.checks)["exit_code_zero"] is True
    assert dict(res.checks)["file_exists"] is True
    assert dict(res.checks)["nonzero_size"] is True
    assert dict(res.checks)["format_valid"] is True


def test_validate_missing_file_failure(tmp_path):
    req = CodeActRequest(goal="summary", output_filename="missing.docx")
    attempt = ExecutionAttempt(
        attempt_number=1,
        code="pass",
        stdout="",
        stderr="",
        traceback=None,
        exit_code=0,
        duration_ms=50,
    )

    res = validate(req, tmp_path, attempt)
    assert res.passed is False
    assert dict(res.checks)["file_exists"] is False


def test_validate_nonzero_exit_code_failure(tmp_path):
    out_file = tmp_path / "out.txt"
    out_file.write_text("hello world " * 20, encoding="utf-8")

    req = CodeActRequest(goal="test", output_filename="out.txt")
    attempt = ExecutionAttempt(
        attempt_number=1,
        code="raise RuntimeError()",
        stdout="",
        stderr="RuntimeError",
        traceback="RuntimeError",
        exit_code=1,
        duration_ms=50,
    )

    res = validate(req, tmp_path, attempt)
    assert res.passed is False
    assert dict(res.checks)["exit_code_zero"] is False


def test_validate_pptx_format(tmp_path):
    from pptx import Presentation

    pptx_file = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Valid Slide"
    prs.save(str(pptx_file))

    req = CodeActRequest(goal="deck", output_filename="deck.pptx")
    attempt = ExecutionAttempt(
        attempt_number=1,
        code="",
        stdout="",
        stderr="",
        traceback=None,
        exit_code=0,
        duration_ms=200,
    )

    res = validate(req, tmp_path, attempt)
    assert res.passed is True
    assert dict(res.checks)["format_valid"] is True


def test_validate_corrupt_pptx_format(tmp_path):
    corrupt_file = tmp_path / "bad.pptx"
    corrupt_file.write_bytes(b"corrupt binary junk " * 100)

    req = CodeActRequest(goal="deck", output_filename="bad.pptx")
    attempt = ExecutionAttempt(
        attempt_number=1,
        code="",
        stdout="",
        stderr="",
        traceback=None,
        exit_code=0,
        duration_ms=100,
    )

    res = validate(req, tmp_path, attempt)
    assert res.passed is False
    assert dict(res.checks)["format_valid"] is False
