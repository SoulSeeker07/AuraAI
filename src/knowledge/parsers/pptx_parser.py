"""
PPTX (PowerPoint) Parser for RAG 2.0 Knowledge Intelligence

Supports:
- Slide-level chunking
- Title extraction
- Text content extraction
"""

import logging
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    logger.warning("python-pptx not installed. PPTX parsing will be limited.")
    Presentation = None

from ..models import DocumentChunk, DocumentMetadata

logger = logging.getLogger(__name__)


class PPTXParser:
    """Parse PPTX files into structured chunks."""

    def __init__(self):
        self.supported_extensions = [".pptx"]

    def supports(self, file_path: Path) -> bool:
        """Check if file type is supported."""
        return file_path.suffix.lower() in self.supported_extensions

    def parse(self, file_path: Path) -> list[DocumentChunk]:
        """Parse PPTX file into slide chunks."""
        chunks = []
        if not Presentation:
            logger.error("python-pptx not installed. Cannot parse PPTX files.")
            return chunks

        try:
            prs = Presentation(file_path)

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []

                # Extract slide title
                slide_title = self._extract_title(slide)

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)

                # Add if there's content
                if slide_content or slide_title:
                    full_content = "\n".join(slide_content)
                    if slide_title:
                        full_content = f"{slide_title}\n\n{full_content}"

                    chunk = DocumentChunk(
                        id=f"{file_path.stem}_slide_{slide_num}",
                        content=full_content,
                        metadata=DocumentMetadata(
                            source=str(file_path),
                            file_type="pptx",
                            chunk_type="slide",
                            chunk_id=f"slide_{slide_num}",
                            line_start=slide_num,
                            line_end=slide_num + 1,
                            slide_number=slide_num,
                            slide_title=slide_title,
                            text_count=len(slide_content),
                        ),
                        embeddings=None,
                    )
                    chunks.append(chunk)

        except Exception as e:
            logger.error(f"Error parsing PPTX file {file_path}: {e}")

        return chunks

    def extract_metadata(self, file_path: Path) -> DocumentMetadata:
        """Extract metadata from PPTX file."""
        if not Presentation:
            return DocumentMetadata(
                source=str(file_path),
                file_type="pptx",
                chunk_type="file",
                chunk_id=file_path.stem,
                line_start=1,
                line_end=0,
            )

        try:
            prs = Presentation(file_path)

            return DocumentMetadata(
                source=str(file_path),
                file_type="pptx",
                chunk_type="file",
                chunk_id=file_path.stem,
                line_start=1,
                line_end=len(prs.slides),
                slide_count=len(prs.slides),
                word_count=sum(
                    sum(
                        len(shape.text.split())
                        for shape in slide.shapes
                        if hasattr(shape, "text")
                    )
                    for slide in prs.slides
                ),
            )
        except Exception as e:
            logger.error(f"Error extracting metadata from PPTX file {file_path}: {e}")
            return DocumentMetadata(
                source=str(file_path),
                file_type="pptx",
                chunk_type="file",
                chunk_id=file_path.stem,
                line_start=1,
                line_end=0,
            )

    def _extract_title(self, slide) -> str | None:
        """Extract slide title."""
        try:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Title is usually on the first shape
                    return shape.text.strip()
            return None
        except Exception:
            return None
